from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define title and content slide layouts
slide_layout_title = prs.slide_layouts[0]  # Title Slide layout
slide_layout_content = prs.slide_layouts[1]  # Content layout

# Helper function to set text style
def set_text_style(text_frame, font_size=24, bold=False, color=(0, 0, 0), alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    text_frame.paragraphs[0].alignment = alignment
    return p

# Slide 1: Title Slide
slide = prs.slides.add_slide(slide_layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Study Smart - The Ultimate Student Companion"
subtitle.text = "Revolutionizing Learning with an All-in-One Toolset\nPresented by [Your Name]"

# Customize title slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Slide 2: Introduction with background style
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
title.text_frame.paragraphs[0].font.bold = True
content.text = (
    "Overview of Study Smart:\n"
    "- All-in-one toolkit tailored for students to enhance productivity.\n"
    "- Empowers students to overcome learning challenges effortlessly.\n"
    "- Offers a variety of features, from mathematical tools to career exploration.\n"
    "- Inspired by the need for smarter and more accessible learning."
)
# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 255)

# Slide 3: Computational Thinking
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Computational Thinking"
content.text = (
    "Why Study Smart fits computational thinking principles?\n"
    "- **Decomposition**: Breaks learning tasks into smaller, manageable parts.\n"
    "- **Pattern Recognition**: Recognizes common learning challenges and offers tools.\n"
    "- **Abstraction**: Simplifies complex learning processes.\n"
    "- **Algorithms**: Implements step-by-step procedures for solving tasks."
)
# Style text
set_text_style(content.text_frame, font_size=22, bold=True, color=(50, 50, 50))

# Slide 4: Features Overview with icon-style bullets
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Features Overview"
content.text = (
    "1. **Graphing Tool**: Plot mathematical functions.\n"
    "2. **Quadratic Equation Solver**: Provides roots and graphs.\n"
    "3. **Equation Solver**: Supports linear and polynomial equations.\n"
    "4. **Unit Converter**: Converts various units of measurement.\n"
    "5. **PDF Reader with TTS**: Reads PDFs aloud for easy listening.\n"
    "6. **Scientific Calculator**: Handles complex mathematical calculations.\n"
    "7. **Mathematical Tools**: Includes geometry, trigonometry, etc.\n"
    "8. **Grade Calculator**: Calculates predicted grades based on scores.\n"
    "9. **Exam Anxiety Remover**: Offers relaxation exercises.\n"
    "10. **Text-to-Handwriting Converter**: Creates handwritten-style notes.\n"
    "11. **Career Exploration**: Suggests careers based on skills.\n"
    "12. **Scholarship Finder**: Helps find available scholarships.\n"
    "13. **To-Do List (Firebase Integration)**: Tracks tasks in real-time."
)
# Style background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(230, 250, 230)

# Slide 5: Firebase Integration with visuals
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Firebase Integration"
content.text = (
    "Why integrate Firebase?\n"
    "- **Real-Time Sync**: Keeps the to-do list updated across devices.\n"
    "- **User Authentication**: Secures user login and data.\n"
    "- **Efficient Data Management**: Stores and retrieves user information quickly.\n"
    "Benefits:\n"
    "- **Smooth Experience**: Real-time updates create a seamless user experience.\n"
    "- **Data Reliability**: Firebase's backend ensures data consistency."
)

# Slide 6: Visual Design for Computational Thinking
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Computational Thinking in Action"
content.text = (
    "- **Decomposition**: Different tools solve specific problems, reducing complexity.\n"
    "- **Pattern Recognition**: Analyzes user input and applies relevant functions.\n"
    "- **Abstraction**: Presents only necessary information to avoid overload.\n"
    "- **Algorithms**: Uses procedures like solving equations and processing conversions."
)

# Slide 7: Use Cases & Target Audience
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Use Cases & Target Audience"
content.text = (
    "Who can benefit from Study Smart?\n"
    "- **Students**: Ideal for exam preparation, project work, and daily tasks.\n"
    "- **Educators**: Provides tools for enhancing classroom engagement.\n"
    "- **Parents**: Helps in guiding children through studies and career choices.\n"
    "Applications:\n"
    "- Homework assistance, exam readiness, and overall study improvement."
)

# Slide 8: Future Enhancements
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Future Enhancements"
content.text = (
    "- **AI-Powered Tutoring**: Personalized learning experience with AI-driven insights.\n"
    "- **Video Tutorials Integration**: Interactive video content for various subjects.\n"
    "- **Cloud-Based Note Storage**: Allows students to store and access notes anywhere."
)

# Slide 9: Closing Slide with Thank You
slide = prs.slides.add_slide(slide_layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You for Your Attention!"
subtitle.text = "For further inquiries, reach out via [Contact Information]\nAny Questions?"

# Save the presentation
prs.save("Enhanced_Study_Smart_Presentation.pptx")

print("Enhanced PowerPoint presentation created successfully as 'Enhanced_Study_Smart_Presentation.pptx'")
