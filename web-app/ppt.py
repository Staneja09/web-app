from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a slide layout (title and content)
slide_layout_title = prs.slide_layouts[0]  # Title Slide layout
slide_layout_content = prs.slide_layouts[1]  # Content layout

# Slide 1: Title Slide
slide = prs.slides.add_slide(slide_layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Study Smart - The Ultimate Student Companion"
subtitle.text = "A Comprehensive Toolset for Smarter Learning\nPresented by [Your Name]"

# Slide 2: Introduction
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Overview of Study Smart:\n"
    "- An all-in-one toolkit for students.\n"
    "- Makes learning more accessible, efficient, and engaging.\n"
    "- Inspired by the need to overcome learning challenges."
)

# Slide 3: Computational Thinking
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Computational Thinking"
content.text = (
    "Why Does Study Smart Lie Under Computational Thinking?\n"
    "- Decomposition: Breaking down learning into smaller tasks.\n"
    "- Pattern Recognition: Solving common study challenges.\n"
    "- Abstraction: Simplifying complex functions.\n"
    "- Algorithms: Step-by-step procedures for various tools."
)

# Slide 4: Features Overview
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Features Overview"
content.text = (
    "- Graphing Tool\n"
    "- Quadratic Equation Solver\n"
    "- Equation Solver\n"
    "- Unit Converter\n"
    "- PDF Reader with TTS (Text-to-Speech)\n"
    "- Scientific Calculator\n"
    "- Mathematical Tools\n"
    "- Grade Calculator\n"
    "- Exam Anxiety Remover Tool\n"
    "- Text-to-Handwriting\n"
    "- Career Exploration Tool\n"
    "- Scholarship Finder Tool\n"
    "- To-Do List with Firebase Integration"
)

# Slide 5: How Each Tool Works (1/2)
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "How Each Tool Works (1/2)"
content.text = (
    "- Graphing Tool: Input equations and visualize functions.\n"
    "- Quadratic Equation Solver: Step-by-step solutions.\n"
    "- Equation Solver: Solves linear, quadratic, and higher-order equations.\n"
    "- Unit Converter: Converts units for various quantities."
)

# Slide 6: How Each Tool Works (2/2)
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "How Each Tool Works (2/2)"
content.text = (
    "- PDF Reader with TTS: Reads PDFs aloud for auditory learning.\n"
    "- Scientific Calculator: Performs complex calculations.\n"
    "- Grade Calculator: Predicts final grades and GPA.\n"
    "- Exam Anxiety Remover: Offers relaxation techniques.\n"
    "- Text-to-Handwriting: Converts typed text to handwriting style.\n"
    "- Career Exploration & Scholarship Finder: Provides career and scholarship information."
)

# Slide 7: Why Does It Work?
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Why Does It Work?"
content.text = (
    "- Solves Real Student Problems\n"
    "- Saves Time & Increases Productivity\n"
    "- Enhances Learning with Interactive Tools\n"
    "- Promotes Healthy Habits (Exam Anxiety Remover)"
)

# Slide 8: Firebase Integration
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Firebase Integration"
content.text = (
    "Why Use Firebase?\n"
    "- Real-Time Data Storage for syncing to-do lists.\n"
    "- User Authentication for secure login.\n"
    "- Database Management for storing user data.\n"
    "Benefits of Firebase Integration:\n"
    "- Smooth and responsive user experience.\n"
    "- Keeps user data synchronized and accessible."
)

# Slide 9: Computational Thinking in Action
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Computational Thinking in Action"
content.text = (
    "- Decomposition: Separate tools for specific tasks.\n"
    "- Pattern Recognition: Addressing common study needs.\n"
    "- Abstraction: Simplifying complex interactions.\n"
    "- Algorithmic Design: Step-by-step processes for calculations."
)

# Slide 10: Use Cases & Target Audience
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Use Cases & Target Audience"
content.text = (
    "Who Can Use Study Smart?\n"
    "- Students: For academic tasks and exam preparation.\n"
    "- Educators: As a recommended study aid.\n"
    "- Parents: To help their children study better.\n"
    "Real-Life Applications:\n"
    "- Homework, exam preparation, and daily study routines."
)

# Slide 11: Future Scope & Enhancements
slide = prs.slides.add_slide(slide_layout_content)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Future Scope & Enhancements"
content.text = (
    "- AI-Based Tutoring: Adding a chatbot for assistance.\n"
    "- More Educational Content: Video tutorials.\n"
    "- Cloud Storage: For storing notes and assignments."
)

# Slide 12: Closing Slide
slide = prs.slides.add_slide(slide_layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You!"
subtitle.text = "Any Questions?\nContact Information: [Your Email or Social Media Handle]"

# Save the presentation
prs.save("Study_Smart_Presentation.pptx")

print("PowerPoint presentation created successfully as 'Study_Smart_Presentation.pptx'")
